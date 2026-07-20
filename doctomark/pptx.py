"""PPTX converter — extracts text, images, tables, charts via python-pptx."""

import html
import re
from pathlib import Path
from typing import Optional

import markdownify
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import BaseConverter, ConversionResult


class PptxConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = (".pptx",)

    def convert(self, input_path: Path) -> ConversionResult:
        prs = pptx.Presentation(str(input_path))
        md_content = ""

        for slide_num, slide in enumerate(prs.slides, 1):
            md_content += f"\n\n<!-- Slide {slide_num} -->\n"

            title_shape = slide.shapes.title

            for shape in slide.shapes:
                # --- Pictures ---
                if self._is_picture(shape):
                    md_content += self._handle_picture(shape)

                # --- Tables ---
                elif self._is_table(shape):
                    md_content += self._handle_table(shape)

                # --- Charts ---
                elif shape.has_chart:
                    md_content += self._handle_chart(shape.chart)

                # --- Text frames ---
                elif shape.has_text_frame:
                    text = shape.text.strip()
                    if not text:
                        continue
                    if shape == title_shape:
                        md_content += f"\n# {text}\n"
                    else:
                        md_content += f"\n{text}\n"

            # --- Speaker notes ---
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    md_content += f"\n> **Notes:** {notes_frame.text.strip()}\n"

        return ConversionResult(markdown=md_content.strip())

    # ─── Helpers ───────────────────────────────────────────────

    @staticmethod
    def _is_picture(shape) -> bool:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            return hasattr(shape, "image")
        return False

    @staticmethod
    def _is_table(shape) -> bool:
        return shape.shape_type == MSO_SHAPE_TYPE.TABLE

    def _handle_picture(self, shape) -> str:
        ext = shape.image.content_type.split("/")[-1]
        blob = shape.image.blob

        # Try to get a description
        alt_text = ""
        try:
            alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
        except Exception:
            pass
        if not alt_text:
            alt_text = shape.name

        safe_name = re.sub(r"\W", "_", alt_text or "picture")
        rel_path = self._save_image(blob, shape.image.content_type)
        if rel_path:
            return f"\n![{alt_text}]({rel_path})\n"
        return ""

    def _handle_table(self, shape) -> str:
        table = shape.table
        if not table.rows:
            return ""

        lines = []
        # Header row
        header = []
        for cell in table.rows[0].cells:
            header.append(cell.text.strip())
        lines.append("| " + " | ".join(header) + " |")
        lines.append("|" + "|".join(["---"] * len(header)) + "|")

        # Data rows (skip header)
        for row in list(table.rows)[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")

        return "\n" + "\n".join(lines) + "\n"

    def _handle_chart(self, chart) -> str:
        md = "\n### 图表"
        if chart.has_title:
            md += f": {chart.chart_title.text_frame.text}"
        md += "\n\n"

        try:
            plot = chart.plots[0]
            categories = [c.label for c in plot.categories]
            series_names = [s.name for s in chart.series]
            md += "| 类别 | " + " | ".join(series_names) + " |\n"
            md += "|" + "|".join(["---"] * (len(series_names) + 1)) + "|\n"

            for idx, cat in enumerate(categories):
                vals = [str(chart.series[s].values[idx]) for s in range(len(series_names))]
                md += "| " + str(cat) + " | " + " | ".join(vals) + " |\n"
        except Exception:
            md += "（无法解析图表数据）\n"

        return md
