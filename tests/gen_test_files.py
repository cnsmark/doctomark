#!/usr/bin/env python3
"""生成测试文件用于验证 doctomark 全流程。

创建带图片的 .docx/.pptx/.xlsx/.pdf 测试文件到 test_files/ 目录。
"""

import io
import zipfile
import os
from pathlib import Path

import mammoth
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage


TEST_DIR = Path(__file__).parent / "test_files"
TEST_DIR.mkdir(exist_ok=True)


def make_test_image(name: str, color: str = "blue") -> bytes:
    """Generate a simple solid-color PNG image and return bytes."""
    img = Image.new("RGB", (200, 150), color=color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    print(f"  ✅ 创建图片: {name} ({len(buf.getvalue())} bytes)")
    return buf.getvalue()


def create_docx():
    """Create a .docx with embedded images."""
    # We'll create it using python-docx via markdown → mammoth? No.
    # Let's create a simple HTML with img and convert via mammoth's reverse? No.
    # Actually let's just create a minimal docx with python-docx directly.
    # But python-docx isn't installed. Let's use mammoth to create from HTML... 
    # No, mammoth is docx→html only.
    # Let's download python-docx.

    try:
        from docx import Document
        from docx.shared import Inches as DocxInches

        img1 = make_test_image("docx_img1.png", "red")
        img2 = make_test_image("docx_img2.png", "green")

        doc = Document()
        doc.add_heading("DOCX 测试文档", level=1)
        doc.add_paragraph("这是一段正文，下面是两张图片：")

        # Embed image 1
        buf = io.BytesIO(img1)
        doc.add_picture(buf, width=DocxInches(2))
        doc.add_paragraph("图1：红色测试图")

        # Embed image 2
        buf = io.BytesIO(img2)
        doc.add_picture(buf, width=DocxInches(2))
        doc.add_paragraph("图2：绿色测试图")

        doc.add_heading("表格测试", level=2)
        table = doc.add_table(rows=3, cols=3)
        table.style = "Light Grid Accent 1"
        headers = ["姓名", "年龄", "城市"]
        data = [
            ["张三", "28", "北京"],
            ["李四", "32", "上海"],
        ]
        for j, h in enumerate(headers):
            table.cell(0, j).text = h
        for i, row_data in enumerate(data, 1):
            for j, val in enumerate(row_data):
                table.cell(i, j).text = val

        path = TEST_DIR / "test.docx"
        doc.save(str(path))
        print(f"✅ 创建 DOCX: {path}")
        return path
    except ImportError:
        print("⚠️  python-docx 未安装，跳过 DOCX 测试文件创建")
        return None


def create_pptx():
    """Create a .pptx with images and text."""
    img1 = make_test_image("pptx_img1.png", "purple")
    img2 = make_test_image("pptx_img2.png", "orange")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PPTX 测试演示"
    slide.placeholders[1].text = "包含图片和表格的测试"

    # Slide 2: Image + text
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    left, top = Inches(1), Inches(1)
    slide.shapes.add_picture(io.BytesIO(img1), left, top, width=Inches(3))

    txBox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    txBox.text = "紫色测试图片"

    # Slide 3: Table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(io.BytesIO(img2), Inches(5), Inches(1), width=Inches(3))

    # Add table
    rows, cols = 3, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(3), Inches(4), Inches(2))
    table = table_shape.table
    for i in range(rows):
        for j in range(cols):
            table.cell(i, j).text = f"({i},{j})"

    # Add speaker notes
    slide.notes_slide.notes_text_frame.text = "这是演讲者备注。"

    path = TEST_DIR / "test.pptx"
    prs.save(str(path))
    print(f"✅ 创建 PPTX: {path}")
    return path


def create_xlsx():
    """Create a .xlsx with embedded images."""
    img = make_test_image("xlsx_img.png", "cyan")

    wb = Workbook()
    ws = wb.active
    ws.title = "数据表"

    # Data
    ws["A1"] = "产品"
    ws["B1"] = "销量"
    ws["C1"] = "单价"
    for i, (name, qty, price) in enumerate([
        ("产品A", 100, 9.9),
        ("产品B", 200, 19.9),
        ("产品C", 150, 15.0),
    ], 2):
        ws[f"A{i}"] = name
        ws[f"B{i}"] = qty
        ws[f"C{i}"] = price

    # Embed image
    img_buf = io.BytesIO(img)
    xl_img = XLImage(img_buf)
    xl_img.width = 200
    xl_img.height = 150
    ws.add_image(xl_img, "E2")

    # Second sheet
    ws2 = wb.create_sheet("摘要")
    ws2["A1"] = "总计"
    ws2["B1"] = 450

    path = TEST_DIR / "test.xlsx"
    wb.save(str(path))
    print(f"✅ 创建 XLSX: {path}")
    return path


def create_pdf():
    """Create a simple PDF with text using Pillow."""
    from PIL import Image, ImageDraw, ImageFont

    # Create a multi-page "text as image" PDF
    pages = []
    for i, text in enumerate([
        "PDF 测试文档 - 第 1 页\n\n这是一个测试 PDF 文件。\n用于验证 PDF 文本提取功能。\n\nHello World!",
        "PDF 测试文档 - 第 2 页\n\n第二页的测试内容。\n\n包含更多文本。",
    ]):
        img = Image.new("RGB", (595, 842), "white")  # A4 at 72 DPI
        draw = ImageDraw.Draw(img)
        for j, line in enumerate(text.split("\n")):
            draw.text((50, 50 + j * 25), line, fill="black")
        pages.append(img)

    # Save as PDF
    path = TEST_DIR / "test.pdf"
    pages[0].save(
        str(path),
        "PDF",
        save_all=True,
        append_images=pages[1:],
    )
    print(f"✅ 创建 PDF: {path}")
    return path


def main():
    print("=" * 50)
    print("创建测试文件...")
    print("=" * 50)

    docx_path = create_docx()
    pptx_path = create_pptx()
    xlsx_path = create_xlsx()
    pdf_path = create_pdf()

    print("\n" + "=" * 50)
    print("创建完成:")
    if docx_path: print(f"  📄 {docx_path}")
    if pptx_path: print(f"  📊 {pptx_path}")
    if xlsx_path: print(f"  📈 {xlsx_path}")
    if pdf_path: print(f"  📕 {pdf_path}")
    print("=" * 50)


if __name__ == "__main__":
    main()
